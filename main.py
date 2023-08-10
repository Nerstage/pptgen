#!/usr/bin/env python3

import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from hymntext import Hymn


def hymn_slide(hymn: Hymn) -> None:
    prs = Presentation("template.pptx")
    prs = append_title_slide(
        prs, hymn.title, hymnal=hymn.hymnal, hymn_index=hymn.hymn_index
    )
    for i in hymn.order:
        if i == "refrain":
            prs = append_word_slide(
                prs, hymn.refrain, hymn.hymnal_code, hymn.hymn_index
            )
        else:
            prs = append_word_slide(
                prs, hymn.verses[i - 1], hymn.hymnal_code, hymn.hymn_index
            )
    prs.save("test.pptx")


def append_title_slide(
    prs: Presentation, title: str, hymnal: str = "", hymn_index: int = 0
) -> Presentation:
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = title
    if hymnal != "":
        if hymn_index != 0:
            subtitle = f"{hymnal} #{hymn_index}"
        else:
            subtitle = hymnal
        slide.placeholders[1].text = subtitle

    return prs


def append_word_slide(
    prs: Presentation,
    words: str,
    hymnal_code: str = None,
    hymnal_index: int = None,
    left: float = 1,
    top: float = 1,
    width: float = 1,
    height: float = 1,
) -> Presentation:
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = words
    p.font.size = Pt(20)
    if hymnal_code and hymnal_index:
        note = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(1), Inches(1))
        note.text = f"{hymnal_code} #{hymnal_index}"
    return prs


def main():
    # test = Hymn("O for a Thousand", "UMH", 116)
    test = Hymn("He Is Exalted", "The Faith We Sing", 2070)
    test.hymnal_code = "TFWS"
    print(test.order)
    hymn_slide(test)


if __name__ == "__main__":
    main()
